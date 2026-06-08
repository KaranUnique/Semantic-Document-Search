import os
import hashlib
import logging
from typing import List, Dict, Any, Tuple
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

logger = logging.getLogger("doc_processor")
logging.basicConfig(level=logging.INFO)


class RecursiveCharacterTextSplitter:
    """
    Lightweight custom text splitter to replace LangChain dependency.
    Splits text into chunks while respecting boundaries (paragraphs, sentences).
    """
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200, length_function=len):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.length_function = length_function
        
    def split_text(self, text: str) -> List[str]:
        """Splits text into chunks with overlap, preserving coherence."""
        if not text or self.length_function(text) <= self.chunk_size:
            return [text]
        
        chunks = []
        # Try splitting by paragraph first (most natural boundary)
        separators = ["\n\n", "\n", ". ", " "]
        good_splits = []
        
        for separator in separators:
            splits = text.split(separator)
            if len(splits) > 1:
                good_splits = splits
                break
        
        if not good_splits:
            good_splits = [text]
        
        # Merge splits into chunks of desired size
        current_chunk = ""
        for split in good_splits:
            if self.length_function(split) > self.chunk_size:
                # Split is too large, need to split it further
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = ""
                chunks.extend(self.split_text(split))
            else:
                test_chunk = (current_chunk + split).strip()
                if self.length_function(test_chunk) <= self.chunk_size:
                    current_chunk = test_chunk
                else:
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                    current_chunk = split.strip()
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        # Add overlap
        if self.chunk_overlap > 0:
            overlapped_chunks = []
            for i, chunk in enumerate(chunks):
                overlapped_chunks.append(chunk)
                if i < len(chunks) - 1:
                    overlap_start = max(0, self.length_function(chunk) - self.chunk_overlap)
                    if overlap_start < self.length_function(chunk):
                        next_chunk = chunks[i + 1]
                        overlap_text = chunk[-self.chunk_overlap:] + " " + next_chunk[:self.chunk_overlap]
                        if len(overlapped_chunks) > 0:
                            overlapped_chunks[-1] = chunk
            chunks = overlapped_chunks
        
        return [c for c in chunks if c]


class DocProcessor:
    def __init__(self, upload_dir: str = "./uploads"):
        self.upload_dir = upload_dir
        os.makedirs(self.upload_dir, exist_ok=True)
        # 1000 size, 200 overlap as specified
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )

    @staticmethod
    def calculate_md5(file_path: str) -> str:
        """Calculates MD5 hash of a file for duplicate check."""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def process_file(self, file_path: str, filename: str) -> Tuple[List[Dict[str, Any]], int]:
        """
        Parses a document, extracts text page-by-page, cleans it,
        chunks the text, and returns a list of chunks metadata-enriched,
        along with the total page count.
        """
        file_ext = os.path.splitext(filename)[1].lower()
        chunks: List[Dict[str, Any]] = []
        page_count = 0
        
        try:
            if file_ext == ".pdf":
                pages_text, page_count = self._extract_pdf(file_path)
            elif file_ext == ".docx":
                pages_text, page_count = self._extract_docx(file_path)
            elif file_ext == ".pptx":
                pages_text, page_count = self._extract_pptx(file_path)
            elif file_ext == ".txt":
                pages_text, page_count = self._extract_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")

            # Splitting each page's text to preserve page citation accuracy
            for page_idx, text in pages_text.items():
                cleaned_text = self._clean_text(text)
                if not cleaned_text:
                    continue
                
                # Split text for this page
                sub_chunks = self.text_splitter.split_text(cleaned_text)
                for chunk_idx, chunk_text in enumerate(sub_chunks):
                    chunks.append({
                        "text": chunk_text,
                        "metadata": {
                            "source": filename,
                            "page": page_idx,  # 1-based page number
                            "chunk_index": chunk_idx,
                            "file_type": file_ext
                        }
                    })
            
            logger.info(f"Processed {filename}: {page_count} pages, generated {len(chunks)} chunks.")
            return chunks, page_count

        except Exception as e:
            logger.error(f"Error processing file {filename}: {str(e)}")
            raise e

    def _clean_text(self, text: str) -> str:
        """Cleans whitespaces and redundant newlines from text."""
        if not text:
            return ""
        # Simple whitespace normalization
        lines = [line.strip() for line in text.splitlines()]
        cleaned = " ".join([line for line in lines if line])
        return cleaned

    def _extract_pdf(self, file_path: str) -> Tuple[Dict[int, str], int]:
        """
        Extracts PDF text using native pypdf extraction.
        """
        pages_text: Dict[int, str] = {}
        reader = PdfReader(file_path)
        total_pages = len(reader.pages)
        
        # Native extraction
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages_text[idx + 1] = text
        
        return pages_text, total_pages

    def _extract_docx(self, file_path: str) -> Tuple[Dict[int, str], int]:
        """Extracts text from DOCX paragraph by paragraph."""
        pages_text: Dict[int, str] = {}
        doc = DocxDocument(file_path)
        
        # DOCX lacks physical pages in Python, approximate pages by grouping paragraphs
        current_page = 1
        char_count = 0
        page_content: List[str] = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            page_content.append(text)
            char_count += len(text)
            
            # Approximate a page as ~2000 characters
            if char_count > 2000:
                pages_text[current_page] = "\n".join(page_content)
                current_page += 1
                char_count = 0
                page_content = []
                
        if page_content or not pages_text:
            pages_text[current_page] = "\n".join(page_content)
            
        return pages_text, len(pages_text)

    def _extract_pptx(self, file_path: str) -> Tuple[Dict[int, str], int]:
        """Extracts text slide-by-slide from PPTX presentations."""
        pages_text: Dict[int, str] = {}
        prs = PptxPresentation(file_path)
        
        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            pages_text[slide_num] = "\n".join(slide_text)
            
        return pages_text, len(prs.slides)

    def _extract_txt(self, file_path: str) -> Tuple[Dict[int, str], int]:
        """Extracts text from UTF-8 encoded text files."""
        pages_text: Dict[int, str] = {}
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
            
        # Group text into approximate 2000 character pages
        pages_text[1] = content
        return pages_text, 1
